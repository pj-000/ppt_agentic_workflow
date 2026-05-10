"""
backend/harness/agents/document_summary.py

封装 document-summarizer skill（vendor/document_summarizer）。
将 PDF/Word 提取的原始文本转换为结构化摘要，供 PPT 生成 pipeline 使用。

流程：
  extract_document_content()     → 提取文本 + 表格
  summarize_document()            → 调用 LLM，按 SKILL.md 规范输出 JSON
  DocumentSummary                  → Pydantic 模型，对应 SKILL.md 的输出 schema
"""
from __future__ import annotations

import json
import logging
import os
import re
import subprocess
import tempfile
import uuid
from pathlib import Path
from typing import Any, Callable, Literal, Optional

from openai import OpenAI
from pydantic import BaseModel, Field, field_validator
import config as pptagent_config
from backend.harness.runtime import (
    HarnessTrace,
    PromptComposer,
    PromptSection,
    RepairOrchestrator,
    SkillContext,
    SkillLoader,
    merge_prompt_sections,
)

logger = logging.getLogger(__name__)

PROJECT_ROOT = Path(__file__).resolve().parents[3]
PROCESSOR_ROOT = PROJECT_ROOT / "vendor" / "document_processor"
_COMPOSER = PromptComposer()


def _record_prompt_bundle(
    *,
    harness_trace: HarnessTrace | None,
    stage: str,
    mode: str,
    context: SkillContext,
    bundle,
    attempt: int | None = None,
    error_signature: str = "",
) -> None:
    if not harness_trace or not bundle.loaded_records:
        return
    harness_trace.record(
        stage=stage,
        payload=bundle.to_trace_payload(
            mode=mode,
            context=context,
            attempt=attempt,
            error_signature=error_signature,
        ),
    )


def skill_paths() -> dict:
    """返回所有 document-processing 相关 skill 的本地路径。"""
    return {
        "document_summarizer": str(PROJECT_ROOT / "backend" / "harness" / "skills" / "document-understanding" / "SKILL.md"),
        "pdf_processor": str(PROCESSOR_ROOT / "pdf" / "SKILL.md"),
        "docx_processor": str(PROCESSOR_ROOT / "docx" / "SKILL.md"),
        "markdown_processor": str(PROCESSOR_ROOT / "markdown" / "SKILL.md"),
        "pptx_processor": str(PROCESSOR_ROOT / "pptx" / "SKILL.md"),
    }


def assert_skills_present() -> None:
    """验证所有必需 skill 文件存在。"""
    for name, path in skill_paths().items():
        if not Path(path).exists():
            raise FileNotFoundError(
                f"Skill 文件缺失: {name} → {path}\n"
                "请确保已将 vendor/document_summarizer 和 vendor/document_processor 目录 clone 到本地。"
            )


# ---------------------------------------------------------------------------
# Pydantic models — mirror the SKILL.md output schema
# ---------------------------------------------------------------------------

class TableSpec(BaseModel):
    description: str
    headers: list[str]
    rows: list[list[str]]


class SectionSummary(BaseModel):
    section_title: str
    section_summary: str
    key_points: list[str] = Field(default_factory=list)
    tables: list[TableSpec] = Field(default_factory=list)
    important_data: list[str] = Field(default_factory=list)
    suggested_slide_count: int = Field(default=1, ge=1, le=10)
    narrative_transition: Optional[str] = None

    @field_validator("key_points", mode="before")
    @classmethod
    def _ensure_list(cls, v: Any) -> list[str]:
        if isinstance(v, list):
            return [str(x) for x in v]
        return []


def _coerce_text(value: Any, default: str = "") -> str:
    if value is None:
        return default
    if isinstance(value, str):
        return value.strip() or default
    if isinstance(value, (int, float, bool)):
        return str(value)
    return default


def _coerce_string_list(value: Any) -> list[str]:
    if isinstance(value, list):
        return [_coerce_text(item) for item in value if _coerce_text(item)]
    if isinstance(value, str):
        parts = re.split(r"[\n；;•·]|(?:\s*[-*]\s+)", value)
        return [part.strip() for part in parts if part and part.strip()]
    return []


def _infer_slide_count(section: dict[str, Any]) -> int:
    direct_value = section.get("suggested_slide_count")
    try:
        if direct_value is not None and str(direct_value).strip():
            return max(1, min(10, int(float(str(direct_value).strip()))))
    except (TypeError, ValueError):
        pass

    complexity = _coerce_text(section.get("complexity")).lower()
    if any(token in complexity for token in ("高", "high", "complex")):
        return 2
    if any(token in complexity for token in ("中-高", "较高", "medium-high")):
        return 2

    key_points = _coerce_string_list(section.get("key_points"))
    important_data = _coerce_string_list(section.get("important_data"))
    summary = _coerce_text(
        section.get("section_summary")
        or section.get("summary")
        or section.get("content")
        or section.get("description")
    )
    title = _coerce_text(section.get("section_title") or section.get("title"))

    evidence_score = 0
    if len(key_points) >= 4:
        evidence_score += 1
    if len(important_data) >= 2:
        evidence_score += 1
    if len(summary) >= 180:
        evidence_score += 1
    if re.search(r"(method|approach|architecture|experiment|result|analysis|framework|算法|方法|架构|实验|结果|分析|系统)", title, re.I):
        evidence_score += 1

    return 2 if evidence_score >= 2 else 1


def _normalize_section_item(item: Any, fallback_title: str = "") -> dict[str, Any] | None:
    if isinstance(item, str):
        title = fallback_title or "未命名章节"
        summary = item.strip()
        if not summary:
            return None
        return {
            "section_title": title,
            "section_summary": summary,
            "key_points": [],
            "tables": [],
            "important_data": [],
            "suggested_slide_count": 1,
            "narrative_transition": None,
        }

    if not isinstance(item, dict):
        return None

    title = _coerce_text(
        item.get("section_title")
        or item.get("title")
        or item.get("heading")
        or item.get("name"),
        fallback_title or "未命名章节",
    )
    summary = _coerce_text(
        item.get("section_summary")
        or item.get("summary")
        or item.get("content")
        or item.get("description")
        or item.get("overview"),
    )
    key_points = _coerce_string_list(item.get("key_points") or item.get("highlights") or item.get("bullet_points"))
    important_data = _coerce_string_list(item.get("important_data") or item.get("metrics") or item.get("numbers"))

    if not summary:
        if key_points:
            summary = "；".join(key_points[:3])
        elif important_data:
            summary = "；".join(important_data[:2])
        else:
            summary = title

    normalized = {
        "section_title": title,
        "section_summary": summary,
        "key_points": key_points,
        "tables": item.get("tables") if isinstance(item.get("tables"), list) else [],
        "important_data": important_data,
        "suggested_slide_count": _infer_slide_count(item),
        "narrative_transition": _coerce_text(
            item.get("narrative_transition") or item.get("transition") or item.get("next_section_hint"),
            None,
        ),
    }
    return normalized


def _normalize_sections_payload(data: dict[str, Any]) -> list[dict[str, Any]]:
    raw_sections = data.get("sections")
    normalized_sections: list[dict[str, Any]] = []

    if isinstance(raw_sections, list):
        for item in raw_sections:
            normalized = _normalize_section_item(item)
            if normalized:
                normalized_sections.append(normalized)
        return normalized_sections

    if not isinstance(raw_sections, dict):
        return normalized_sections

    section_order = data.get("section_order")
    ordered_keys = [str(item) for item in section_order] if isinstance(section_order, list) else list(raw_sections.keys())
    title_value = _coerce_text(raw_sections.get("Title"))
    abstract_value = _coerce_text(raw_sections.get("Abstract"))
    authors_value = _coerce_text(raw_sections.get("Authors"))

    if not data.get("document_title") and title_value:
        data["document_title"] = title_value
    if not data.get("document_summary"):
        summary_parts = [part for part in (abstract_value, authors_value and f"作者：{authors_value}") if part]
        if summary_parts:
            data["document_summary"] = " ".join(summary_parts[:2])

    skip_keys = {"title", "authors"}
    for key in ordered_keys:
        value = raw_sections.get(key)
        if value is None:
            continue
        lowered = str(key).strip().lower()
        if lowered in skip_keys:
            continue
        if lowered == "abstract":
            if not data.get("document_summary") and isinstance(value, str) and value.strip():
                data["document_summary"] = value.strip()
            continue
        normalized = _normalize_section_item(value, fallback_title=str(key))
        if normalized:
            normalized_sections.append(normalized)

    return normalized_sections


def _normalize_document_summary_payload(data: dict[str, Any]) -> dict[str, Any]:
    normalized = dict(data)
    normalized["sections"] = _normalize_sections_payload(normalized)

    if not normalized.get("document_title"):
        first_title = normalized["sections"][0]["section_title"] if normalized["sections"] else "未命名文档"
        normalized["document_title"] = first_title

    if not normalized.get("document_summary"):
        section_summaries = [item["section_summary"] for item in normalized["sections"][:2] if item.get("section_summary")]
        normalized["document_summary"] = " ".join(section_summaries)[:500] if section_summaries else normalized["document_title"]

    metadata = normalized.get("metadata")
    normalized["metadata"] = metadata if isinstance(metadata, dict) else {}
    if "page_count" in normalized["metadata"] and "total_pages" not in normalized["metadata"]:
        normalized["metadata"]["total_pages"] = normalized["metadata"].pop("page_count")

    hints = normalized.get("ppt_generation_hints")
    normalized["ppt_generation_hints"] = hints if isinstance(hints, dict) else {}
    if not normalized["ppt_generation_hints"].get("suggested_total_slides"):
        body = sum(int(item.get("suggested_slide_count") or 1) for item in normalized["sections"])
        normalized["ppt_generation_hints"]["suggested_total_slides"] = max(
            4,
            min(pptagent_config.MAX_PPT_SLIDES, body + 2),
        )

    return normalized


class Metadata(BaseModel):
    source_file: Optional[str] = None
    total_pages: Optional[int] = None
    language: str = "auto-detect"
    estimated_reading_time_minutes: Optional[int] = None


class PPTGenerationHints(BaseModel):
    suggested_total_slides: int = Field(ge=4, le=pptagent_config.MAX_PPT_SLIDES, default=8)
    audience: str = "general"
    style_preference: str = "informative"
    recommended_visual_elements: list[str] = Field(default_factory=list)
    content_focus: str = "knowledge_sharing"
    special_requirements: Optional[str] = None


class DocumentSummary(BaseModel):
    """对应 SKILL.md 的完整输出 schema。"""
    document_title: str
    document_summary: str
    sections: list[SectionSummary] = Field(min_length=1)
    metadata: Metadata = Field(default_factory=Metadata)
    ppt_generation_hints: PPTGenerationHints = Field(default_factory=PPTGenerationHints)

    # 原始输入的影子副本，供调试用
    _raw_text_preview: Optional[str] = None
    _source_file: Optional[str] = None

    def total_suggested_slides(self) -> int:
        """封面 + 各 section slides + closing。"""
        body = sum(s.suggested_slide_count for s in self.sections)
        return min(body + 2, pptagent_config.MAX_PPT_SLIDES)

    def to_planner_context(self) -> str:
        """
        将摘要转换为 planner 可直接消费的字符串格式。
        用于注入到 PPT outline planning 的 user prompt 中。
        """
        lines = [
            _COMPOSER.load_document_planner_context_header_template().format(
                document_title=self.document_title,
                document_summary=self.document_summary,
            ).rstrip(),
        ]
        for i, section in enumerate(self.sections):
            key_points_block = ""
            if section.key_points:
                key_points_block = "**要点**:\n" + "\n".join(f"  - {point}" for point in section.key_points) + "\n"

            important_data_block = ""
            if section.important_data:
                important_data_block = "**关键数据**:\n" + "\n".join(f"  - {data}" for data in section.important_data) + "\n"

            tables_block_parts: list[str] = []
            for table in section.tables:
                table_lines: list[str] = []
                if table.headers:
                    table_lines.append("  | " + " | ".join(table.headers) + " |")
                    table_lines.append("  | " + " | ".join(["---"] * len(table.headers)) + " |")
                for row in table.rows[:5]:
                    table_lines.append("  | " + " | ".join(str(c) for c in row) + " |")
                if len(table.rows) > 5:
                    table_lines.append("  | ... |")
                tables_block_parts.append(
                    _COMPOSER.load_document_planner_context_table_template().format(
                        table_description=table.description,
                        table_body="\n".join(table_lines),
                    ).rstrip()
                )
            tables_block = ("\n".join(tables_block_parts) + "\n") if tables_block_parts else ""

            lines.append(
                _COMPOSER.load_document_planner_context_section_template().format(
                    section_index=i + 1,
                    section_title=section.section_title,
                    section_summary=section.section_summary,
                    key_points_block=key_points_block,
                    important_data_block=important_data_block,
                    tables_block=tables_block,
                    narrative_transition_line=(
                        f"*→ 下一节: {section.narrative_transition}*\n"
                        if section.narrative_transition
                        else ""
                    ),
                ).rstrip()
            )

        hints = self.ppt_generation_hints
        lines.append(
            _COMPOSER.load_document_planner_context_hints_template().format(
                suggested_total_slides=hints.suggested_total_slides,
                audience=hints.audience,
                style_preference=hints.style_preference,
                recommended_visual_elements_line=(
                    f"- 推荐视觉元素: {', '.join(hints.recommended_visual_elements)}\n"
                    if hints.recommended_visual_elements
                    else ""
                ),
                content_focus_line=(f"- 内容类型: {hints.content_focus}\n" if hints.content_focus else ""),
                special_requirements_line=(
                    f"- 特殊要求: {hints.special_requirements}\n"
                    if hints.special_requirements
                    else ""
                ),
            ).rstrip()
        )

        return "\n".join(lines)


# ---------------------------------------------------------------------------
# Document extraction (PDF / DOCX → raw text + tables)
# ---------------------------------------------------------------------------

def extract_text_from_pdf(file_path: str) -> tuple[str, list[list[list[str]]], int]:
    """
    提取 PDF 文本和表格。
    优先使用 pdfplumber（保留布局），fallback 到 pypdf。

    Returns:
        (extracted_text, tables, page_count)
    """
    try:
        import pdfplumber

        text_parts: list[str] = []
        all_tables: list[list[list[str]]] = []
        page_count = 0

        with pdfplumber.open(file_path) as pdf:
            page_count = len(pdf.pages)
            for page in pdf.pages:
                t = page.extract_text()
                if t:
                    text_parts.append(t.strip())

                page_tables = page.extract_tables()
                for pt in page_tables:
                    if pt and len(pt) >= 2:
                        all_tables.append(pt)

        return "\n\n".join(text_parts), all_tables, page_count

    except ImportError:
        logger.warning("[DocumentSummarizer] pdfplumber 不可用，fallback 到 pypdf")
    except Exception as exc:
        logger.warning(f"[DocumentSummarizer] pdfplumber 提取失败: {exc}")

    # Fallback: pypdf (no table extraction)
    try:
        from pypdf import PdfReader
        reader = PdfReader(file_path)
        page_count = len(reader.pages)
        text_parts = []
        for page in reader.pages:
            t = page.extract_text()
            if t:
                text_parts.append(t.strip())
        return "\n\n".join(text_parts), [], page_count
    except Exception as exc:
        logger.error(f"[DocumentSummarizer] pypdf 提取也失败: {exc}")
        return "", [], 0


def extract_text_from_docx(file_path: str) -> tuple[str, list[list[list[str]]], int]:
    """
    提取 DOCX 文本。表格用 python-docx 提取。

    Returns:
        (extracted_text, tables, page_count_estimate)
    """
    try:
        from docx import Document

        doc = Document(file_path)
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]

        # 提取表格
        tables: list[list[list[str]]] = []
        for table in doc.tables:
            rows_data: list[list[str]] = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows_data.append(cells)
            if rows_data:
                tables.append(rows_data)

        # DOCX 只能估算页数
        estimated_pages = max(1, len(paragraphs) // 30)

        return "\n\n".join(paragraphs), tables, estimated_pages

    except ImportError:
        logger.warning("[DocumentSummarizer] python-docx 不可用，尝试 pandoc")
    except Exception as exc:
        logger.warning(f"[DocumentSummarizer] python-docx 提取失败: {exc}")

    # Fallback: pandoc
    try:
        result = subprocess.run(
            ["pandoc", file_path, "-o", "-", "--to", "plain"],
            capture_output=True, text=True, timeout=60,
        )
        if result.returncode == 0:
            return result.stdout, [], 0
    except Exception as exc:
        logger.warning(f"[DocumentSummarizer] pandoc 提取也失败: {exc}")

    return "", [], 0


def _strip_yaml_frontmatter(text: str) -> str:
    """Strip YAML frontmatter (content between --- fences at the top)."""
    if text.startswith("---"):
        end = text.find("\n---", 3)
        if end != -1:
            return text[end + 5:].lstrip()
    return text


def _extract_markdown_tables(text: str) -> list[list[list[str]]]:
    """Extract Markdown tables (| Header | ... | style) as 2D arrays."""
    import re

    tables: list[list[list[str]]] = []
    lines = text.split("\n")
    current_table: list[list[str]] = []

    for line in lines:
        line = line.strip()
        if not line.startswith("|"):
            if current_table:
                tables.append(current_table)
                current_table = []
            continue

        if re.match(r"^\|[\s\-:|]+\|$", line):
            continue

        cells = [cell.strip() for cell in line.strip("|").split("|")]
        if cells:
            current_table.append(cells)

    if current_table:
        tables.append(current_table)

    return tables


def extract_text_from_markdown(file_path: str) -> tuple[str, list[list[list[str]]], int]:
    """
    提取 Markdown 文本内容。参考 vendor/document_processor/markdown/SKILL.md。

    Returns:
        (extracted_text, tables, page_count_estimate)
    """
    text = ""
    for enc in ("utf-8", "utf-8-sig", "gbk"):
        try:
            with open(file_path, "r", encoding=enc) as f:
                text = f.read()
            break
        except (UnicodeDecodeError, LookupError):
            continue

    if not text:
        logger.error("[DocumentSummarizer] Markdown 读取失败（所有编码均失败）")
        return "", [], 0

    # Normalize line endings and strip YAML frontmatter
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = _strip_yaml_frontmatter(text)

    tables = _extract_markdown_tables(text)
    estimated_pages = max(1, len(text) // 3000)

    return text.strip(), tables, estimated_pages


def extract_text_from_pptx(file_path: str) -> tuple[str, list[list[list[str]]], int]:
    """
    提取 PPTX 文本。优先使用 python-pptx。

    Returns:
        (extracted_text, tables, page_count_estimate)
    """
    try:
        from pptx import Presentation

        prs = Presentation(file_path)
        text_parts: list[str] = []
        tables: list[list[list[str]]] = []
        slide_count = len(prs.slides)

        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        line = paragraph.text.strip()
                        if line:
                            text_parts.append(line)
                if shape.has_table:
                    table_data: list[list[str]] = []
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        table_data.append(cells)
                    if table_data:
                        tables.append(table_data)

        # PPTX 估算页数
        estimated_pages = max(1, slide_count)
        return "\n\n".join(text_parts), tables, estimated_pages

    except ImportError:
        logger.warning("[DocumentSummarizer] python-pptx 不可用")
    except Exception as exc:
        logger.warning(f"[DocumentSummarizer] python-pptx 提取失败: {exc}")

    # Fallback: markitdown
    try:
        import subprocess

        result = subprocess.run(
            ["python", "-m", "markitdown", file_path],
            capture_output=True, text=True, timeout=60,
        )
        if result.returncode == 0 and result.stdout:
            return result.stdout.strip(), [], 0
    except Exception as exc:
        logger.warning(f"[DocumentSummarizer] markitdown fallback 也失败: {exc}")

    return "", [], 0


def extract_document_content(
    file_path: str,
) -> tuple[str, list[list[list[str]]], int]:
    """
    统一入口：自动根据文件扩展名选择提取方式。

    Returns:
        (raw_text, tables, page_count)
    """
    suffix = Path(file_path).suffix.lower()
    if suffix == ".pdf":
        return extract_text_from_pdf(file_path)
    if suffix in {".docx", ".doc"}:
        return extract_text_from_docx(file_path)
    if suffix == ".md":
        return extract_text_from_markdown(file_path)
    if suffix == ".pptx":
        return extract_text_from_pptx(file_path)

    raise ValueError(f"Unsupported file format: {suffix}")


# ---------------------------------------------------------------------------
# LLM summarization
# ---------------------------------------------------------------------------

def _load_skill() -> str:
    composer = PromptComposer()
    loader = SkillLoader()
    body = loader.read_skill_body("document-understanding").strip()
    repair_patterns = loader.read_reference("document-understanding", "repair_patterns.md").strip()
    return composer.load_document_summary_system_wrapper_template().format(
        skill_body=body,
        repair_patterns=repair_patterns,
    ).strip()


def _build_messages(
    raw_text: str,
    tables: list[list[list[str]]],
    metadata: dict[str, Any],
    skill_md: str,
) -> list[dict[str, str]]:
    """
    构造 LLM messages。
    system: SKILL.md 内容
    user: 原始文本 + 表格 + 元数据
    """
    user_parts: list[str] = []
    if tables:
        user_parts.append(_COMPOSER.load_document_summary_tables_header_template().strip())
        for i, table in enumerate(tables):
            if not table or len(table) < 2:
                continue
            user_parts.append(
                _COMPOSER.load_document_summary_table_heading_template().format(table_index=i + 1).strip()
            )
            for row in table[:10]:  # 最多10行，防止 token 爆炸
                user_parts.append(" | ".join(str(c) for c in row))
            if len(table) > 10:
                user_parts.append(_COMPOSER.load_document_summary_table_truncation_template().strip())
            user_parts.append("")

    prompt = _COMPOSER.load_document_summary_user_prompt_template()
    user_content = prompt
    replacements = {
        "{source_file_line}": (
            _COMPOSER.load_document_summary_source_file_line_template().format(source_file=metadata["source_file"]).strip()
            if metadata.get("source_file") else ""
        ),
        "{page_count_line}": (
            _COMPOSER.load_document_summary_page_count_line_template().format(page_count=metadata["page_count"]).strip()
            if metadata.get("page_count") else ""
        ),
        "{raw_text}": raw_text,
        "{tables_section}": "\n".join(user_parts).strip(),
    }
    for key, value in replacements.items():
        user_content = user_content.replace(key, value)

    return [
        {"role": "system", "content": skill_md},
        {"role": "user", "content": user_content},
    ]


def _parse_summary(raw_content: str) -> DocumentSummary:
    """
    从 LLM 原始响应中解析出 DocumentSummary。
    包含 JSON 修复和解析逻辑，与 planner.py 的 _extract_json 类似。
    """
    cleaned = str(raw_content or "").strip()
    cleaned = re.sub(r"^```(?:json)?\s*", "", cleaned)
    cleaned = re.sub(r"\s*```$", "", cleaned)

    # 提取 JSON blob
    positions = [(cleaned.find("{"), "{"), (cleaned.find("["), "[")]
    positions = [(pos, ch) for pos, ch in positions if pos >= 0]
    if positions:
        start, opener = min(positions, key=lambda item: item[0])
        blob = _extract_json_blob(cleaned[start:])
        if blob:
            cleaned = blob

    # 尝试直接解析
    try:
        data = json.loads(cleaned, strict=False)
    except json.JSONDecodeError:
        # 修复 JSON 后重试
        repaired = _repair_json(cleaned)
        try:
            data = json.loads(repaired, strict=False)
        except json.JSONDecodeError:
            logger.error(f"[DocumentSummarizer] JSON 解析失败，原始内容前300字：{cleaned[:300]}")
            raise ValueError(f"无法解析 LLM 返回的 JSON，内容前300字：{cleaned[:300]}")

    normalized = _normalize_document_summary_payload(data)
    return DocumentSummary.model_validate(normalized)


def _extract_json_blob(text: str) -> str:
    """提取完整的最外层 JSON 对象或数组。"""
    if not text or text[0] not in "{[":
        return text

    opener = text[0]
    closer = "}" if opener == "{" else "]"
    stack = [closer]
    in_string = False
    escape = False

    for i in range(1, len(text)):
        ch = text[i]
        if escape:
            escape = False
            continue
        if ch == "\\":
            escape = True
            continue
        if ch == '"':
            in_string = not in_string
            continue
        if in_string:
            continue
        if ch == closer and stack and ch == stack[-1]:
            stack.pop()
            if not stack:
                return text[:i + 1]
        elif ch == opener:
            stack.append(closer)

    return text


def _repair_json(text: str) -> str:
    """简单的 JSON 修复：处理缺失逗号和未闭合的括号。"""
    # 移除尾随逗号
    text = re.sub(r",(\s*[}\]])", r"\1", text)
    # 补充闭合括号
    result = list(text)
    stack: list[str] = []
    in_string = False
    escape = False

    for ch in text:
        if escape:
            escape = False
            continue
        if ch == "\\":
            escape = True
            continue
        if ch == '"':
            in_string = not in_string
            continue
        if not in_string:
            if ch in "{[":
                stack.append("}" if ch == "{" else "]")
            elif ch in "}]":
                if stack and stack[-1] == ch:
                    stack.pop()

    while stack and stack[-1] in "}])":
        result.append(stack.pop())

    return "".join(result)


def summarize_document(
    raw_text: str,
    *,
    tables: Optional[list[list[list[str]]]] = None,
    source_file: Optional[str] = None,
    page_count: Optional[int] = None,
    language: str = "zh-CN",
    model_provider: str = "minmax",
    thinking_callback: Optional[Callable[[str], None]] = None,
    max_input_chars: int = 80000,
    harness_trace: HarnessTrace | None = None,
) -> DocumentSummary:
    """
    主入口：调用 LLM 将原始文档内容转换为结构化摘要。

    三个 skill 的协作关系：
      pdf / docx SKILL  → 指导如何从文件提取文本和表格
      document-summarizer SKILL → 指导 LLM 如何把原始文本转换为结构化摘要

    Args:
        raw_text:          从 PDF/DOCX 提取的原始文本
        tables:            从 PDF 提取的表格列表
        source_file:       来源文件名
        page_count:        文档页数
        language:          文档语言
        model_provider:    LLM provider，传递给 config
        thinking_callback: 推理过程回调（可选）
        max_input_chars:   截断输入文本的最大字符数（控制 token 消耗）

    Returns:
        DocumentSummary Pydantic 模型

    Raises:
        ValueError: LLM 返回非 JSON 或 JSON 格式错误
        RuntimeError: LLM API 调用失败
    """
    if not raw_text or not raw_text.strip():
        raise ValueError("原始文本为空，请先提取文档内容")

    # 截断超长文本
    truncated = False
    if len(raw_text) > max_input_chars:
        raw_text = raw_text[:max_input_chars]
        truncated = True

    skill_md = _load_skill()

    provider_settings = pptagent_config.get_llm_provider_settings(model_provider)
    client = OpenAI(
        api_key=provider_settings["api_key"],
        base_url=provider_settings["base_url"],
    )
    model_id = provider_settings["model_id"]
    composer = PromptComposer()
    runtime = composer.runtime
    repair_orchestrator = RepairOrchestrator(
        runtime,
        run_id=uuid.uuid4().hex[:8],
        phase="document-understanding",
    )

    metadata: dict[str, Any] = {
        "source_file": source_file,
        "page_count": page_count,
    }
    tables = tables or []

    messages = _build_messages(raw_text, tables, metadata, skill_md)
    trigger_stage = "document_summary"
    layout_scope = "document_batch"
    visual_mode_scope = "summary_json"
    prompt_context = SkillContext(
        phase="document-understanding",
        trigger_stage=trigger_stage,
        layout_scope=layout_scope,
        visual_mode_scope=visual_mode_scope,
        course_type="document",
        provider=model_id,
        language=language,
    )
    prevention_bundle = runtime.build_prevention_bundle(
        context=prompt_context,
        heading="## 长期技能目录（文档理解）",
        max_items=2,
    )
    _record_prompt_bundle(
        harness_trace=harness_trace,
        stage=trigger_stage,
        mode="prevention",
        context=prompt_context,
        bundle=prevention_bundle,
    )

    if truncated:
        messages[1] = {
            "role": "user",
            "content": merge_prompt_sections(
                PromptSection(source_type="static_prompt", identifier="document_summary:user", content=messages[1]["content"]),
                PromptSection(
                    source_type="fallback_notice",
                    identifier="document_summary:truncation_notice",
                    content=composer.load_document_summary_truncation_notice_template().replace(
                        "{max_input_chars}",
                        str(max_input_chars),
                    ),
                ),
            ),
        }
    messages[1] = {
        "role": "user",
        "content": merge_prompt_sections(messages[1]["content"], prevention_bundle),
    }

    from backend.tools.openai_compat import build_chat_completion_kwargs, stream_chat_completion_text

    raw_content = ""
    reasoning = ""
    last_error = ""
    last_error_signature: str | None = None
    for attempt in range(1, 4):
        attempt_messages = list(messages)
        loaded_repair_memory_ids: list[str] = []
        if last_error_signature:
            repair_bundle = runtime.build_repair_bundle(
                context=prompt_context,
                error_signature=last_error_signature,
                max_items=1,
            )
            loaded_repair_memory_ids = list(repair_bundle.runtime_memory_ids)
            _record_prompt_bundle(
                harness_trace=harness_trace,
                stage=trigger_stage,
                mode="repair",
                context=prompt_context,
                bundle=repair_bundle,
                attempt=attempt,
                error_signature=last_error_signature,
            )
            attempt_messages[1] = {
                "role": "user",
                "content": merge_prompt_sections(
                    attempt_messages[1]["content"],
                    repair_bundle,
                    PromptSection(
                        source_type="repair_feedback",
                        identifier="document_summary:retry_feedback",
                        content=composer.load_document_summary_retry_feedback_template().replace(
                            "{retry_feedback_block}",
                            "\n".join(
                                repair_orchestrator.build_retry_feedback(
                                    error=last_error,
                                    error_signature=last_error_signature,
                                    layout_scope=layout_scope,
                                    visual_mode_scope=visual_mode_scope,
                                )
                            ),
                        ),
                    ),
                ),
            }

        raw_content, reasoning = stream_chat_completion_text(
            client,
            model=model_id,
            max_tokens=16384,
            messages=attempt_messages,
            on_reasoning_chunk=thinking_callback,
            **build_chat_completion_kwargs(model_id),
        )

        try:
            summary = _parse_summary(raw_content)
            if last_error_signature:
                repair_instruction = repair_orchestrator.build_repair_instruction(
                    error_signature=last_error_signature,
                    error=last_error,
                    layout_scope=layout_scope,
                    visual_mode_scope=visual_mode_scope,
                )
                repair_orchestrator.remember_success(
                    trigger_stage=trigger_stage,
                    error_signature=last_error_signature,
                    error=last_error,
                    repair_instruction=repair_instruction,
                    layout_scope=layout_scope,
                    visual_mode_scope=visual_mode_scope,
                    course_type_scope="document",
                    provider_scope=model_id,
                    language_scope=language,
                    before_pattern=raw_content[:400],
                    after_pattern=summary.model_dump_json()[:400],
                    conditions=[f"source={Path(source_file).suffix.lower() if source_file else 'unknown'}"],
                )
            break
        except Exception as exc:
            for memory_id in dict.fromkeys(loaded_repair_memory_ids):
                repair_orchestrator.mark_memory_failure(memory_id)
            last_error = str(exc)
            last_error_signature = repair_orchestrator.classify_error(
                last_error,
                stage=trigger_stage,
            )
            if attempt == 3:
                raise

    logger.debug(f"[DocumentSummarizer] LLM reasoning 前100字：{reasoning[:100]}")
    logger.info(
        f"[DocumentSummarizer] 摘要生成完成: "
        f"title={summary.document_title!r}, "
        f"sections={len(summary.sections)}, "
        f"suggested_slides={summary.ppt_generation_hints.suggested_total_slides}"
    )
    return summary
